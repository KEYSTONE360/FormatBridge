from __future__ import annotations

import os
import subprocess
import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from formatbridge.converter import Converter, ConvertOptions, verify_output
from formatbridge.engines import EngineManager


@unittest.skipUnless(os.name == "nt", "Windows COM integration only")
class OfficeIntegrationTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.engines = EngineManager()
        if not cls.engines.ms_office:
            raise unittest.SkipTest("Microsoft Office not detected")

    def setUp(self):
        self.temp = tempfile.TemporaryDirectory(prefix="formatbridge_office_")
        self.root = Path(self.temp.name)
        self.converter = Converter(self.engines)
        self.options = ConvertOptions(quality=85, dpi=120)

    def tearDown(self):
        self.temp.cleanup()

    def test_docx_xlsx_pptx_to_pdf(self):
        docx = self.root / "교차검증.docx"
        document = Document()
        document.add_heading("포맷브릿지 문서 검증", 0)
        document.add_paragraph("한글 DOCX 문서를 PDF로 변환합니다.")
        document.save(docx)

        xlsx = self.root / "교차검증.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "검증"
        sheet.append(["항목", "값"])
        sheet.append(["한글", 2026])
        workbook.save(xlsx)

        pptx = self.root / "교차검증.pptx"
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = "포맷브릿지 프레젠테이션 검증"
        deck.save(pptx)

        for source in (docx, xlsx, pptx):
            with self.subTest(source=source.suffix):
                output = self.converter.convert(source, ".pdf", self.root / "pdf", self.options)[0]
                ok, detail = verify_output(output, ".pdf")
                self.assertTrue(ok, detail)
        docx_pdf = self.root / "pdf" / "교차검증.pdf"
        roundtrip = self.converter.convert(docx_pdf, ".docx", self.root / "roundtrip", self.options)[0]
        ok, detail = verify_output(roundtrip, ".docx")
        self.assertTrue(ok, detail)

    def test_hwp_roundtrip_via_hancom(self):
        if not self.engines.hancom_exe:
            self.skipTest("Hancom Office not detected")
        hwp = self.root / "한글변환.hwp"
        creator = Path(__file__).with_name("create_hwp_sample.ps1")
        result = subprocess.run(
            [
                self.engines.ps32,
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(creator),
                "-OutputPath",
                str(hwp),
            ],
            capture_output=True,
            timeout=240,
            check=False,
        )
        self.assertEqual(result.returncode, 0, (result.stderr or result.stdout).decode(errors="replace"))
        self.assertTrue(hwp.exists() and hwp.stat().st_size > 0)
        with hwp.open("rb") as handle:
            self.assertEqual(handle.read(8), bytes.fromhex("D0CF11E0A1B11AE1"))
        hwpx = None
        for ext in (".pdf", ".hwpx", ".docx", ".txt"):
            with self.subTest(target=ext):
                output = self.root / f"한글변환{ext}"
                self.engines.office_convert("hwp", hwp, output, timeout=240)
                ok, detail = verify_output(output, ext)
                self.assertTrue(ok, detail)
                if ext == ".hwpx":
                    hwpx = output
        self.assertIsNotNone(hwpx)
        hwpx_pdf = self.root / "HWPX재입력.pdf"
        self.engines.office_convert("hwp", hwpx, hwpx_pdf, timeout=240)
        ok, detail = verify_output(hwpx_pdf, ".pdf")
        self.assertTrue(ok, detail)


if __name__ == "__main__":
    unittest.main(verbosity=2)
